from flask import Flask, request, jsonify, render_template, send_file # type: ignore
from flask_cors import CORS # type: ignore
import os
import sys
import webbrowser
import subprocess
import time
import threading
from datetime import datetime
from pptx import Presentation # type: ignore
from pptx.util import Inches, Pt # type: ignore
from pptx.enum.text import PP_ALIGN # type: ignore
from pptx.dml.color import RGBColor # type: ignore
import requests

from ollama_brain import ask_ollama
from executor import execute_plan

# Determine if running as PyInstaller bundle
is_bundled = getattr(sys, 'frozen', False)

# Ollama management
def is_ollama_running():
    """Check if Ollama is running"""
    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=3)
        return response.status_code == 200
    except Exception as e:
        return False

def start_ollama_if_needed():
    """Start Ollama if not running"""
    print("\n🤖 Checking Ollama status...")
    
    # Check if already running
    for attempt in range(3):
        if is_ollama_running():
            print("✅ Ollama is running at http://localhost:11434")
            return True
        if attempt < 2:
            print(f"⏳ Waiting for Ollama... ({attempt + 1}/3)")
            time.sleep(2)
    
    print("⚠️  Ollama is not running")
    print("📝 Please make sure to:")
    print("   1. Run SETUP_OLLAMA.bat to install Ollama")
    print("   2. Keep the Ollama window open while using the app")
    print("   3. Ollama needs to be running at http://localhost:11434")
    
    return False

app = Flask(__name__, static_folder='static', static_url_path='', template_folder='templates')
CORS(app)

PPT_DIR = "generated_ppts"
os.makedirs(PPT_DIR, exist_ok=True)

print("🔥 Zen AutoGPT Agent starting...")
print(f"📁 Static folder: {app.static_folder}")

# Start Ollama on app startup
print("\n🤖 Checking Ollama status...")
start_ollama_if_needed()

@app.route("/ollama-status", methods=["GET"])
def ollama_status():
    """Check Ollama status"""
    if is_ollama_running():
        return jsonify({"status": "running", "url": "http://localhost:11434"})
    else:
        return jsonify({
            "status": "not_running",
            "message": "Ollama is not running. Please run SETUP_OLLAMA.bat first.",
            "setup_instructions": "Run SETUP_OLLAMA.bat to install and start Ollama"
        }), 503

@app.route("/")
def home():
    """Serve the frontend app"""
    try:
        return render_template("index.html")
    except:
        # Fallback for desktop app
        return app.send_static_file("index.html")
@app.route("/agent", methods=["POST"])
def agent():
    # Check Ollama status first
    if not is_ollama_running():
        return jsonify({
            "error": "Ollama is not running",
            "message": "Please start Ollama by running SETUP_OLLAMA.bat first",
            "instructions": "1. Double-click SETUP_OLLAMA.bat\n2. Wait for it to complete\n3. Keep the Ollama window open\n4. Refresh this page"
        }), 503
    
    user_text = request.json.get("text", "")
    user_text_lower = user_text.lower()

    print(f"🧠 User said: {user_text}")

    try:
        plan = ask_ollama(user_text)
        print("📜 Raw Plan:", plan)

        # Simply pass through all actions
        # The executor will handle open_application with fallback to website
        # And handle open_website directly
        
        print("🧹 Final Plan:", plan)

        result = execute_plan(plan)

        return jsonify({
            "reply": result
        })
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({
            "error": f"Error processing request: {str(e)}",
            "message": "Make sure Ollama is running at http://localhost:11434"
        }), 500

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    """Generate PowerPoint presentation from slides with images"""
    try:
        data = request.json
        title = data.get("title", "Generated Presentation")
        slides_data = data.get("slides", [])
        
        if not slides_data:
            return jsonify({"success": False, "error": "No slides provided"})
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Add content slides
        for idx, slide_data in enumerate(slides_data, 1):
            # Use blank layout
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(30, 30, 46)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = slide_data.get("title", f"Slide {idx}")
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 212, 255)
            
            # Add image if available with better sizing
            image_url = slide_data.get("image")
            if image_url:
                try:
                    # Download image
                    img_response = requests.get(image_url, timeout=5)
                    if img_response.status_code == 200:
                        img_path = os.path.join(PPT_DIR, f"temp_img_{idx}.jpg")
                        with open(img_path, "wb") as f:
                            f.write(img_response.content)
                        
                        # Add to slide - better positioning
                        try:
                            # Center the image better
                            slide.shapes.add_picture(img_path, Inches(2), Inches(1.2), width=Inches(6), height=Inches(4))
                            print(f"✅ Image added to slide {idx}")
                        except Exception as e:
                            print(f"⚠️ Could not add image to slide {idx}: {e}")
                except Exception as e:
                    print(f"⚠️ Image download failed for slide {idx}: {e}")
            
            # Add content text with better formatting
            content_text = slide_data.get("content", "")
            
            # Clean up the content text
            if len(content_text) > 200:
                # For longer text, truncate to first sentence or 200 chars
                sentences = content_text.split('.')
                content_text = sentences[0] + '.' if sentences[0] else content_text[:200]
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.8))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = content_text.strip()
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(224, 224, 224)
            p.level = 0
            p.line_spacing = 1.3
        
        # Save presentation
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(PPT_DIR, filename)
        prs.save(filepath)
        
        print(f"✅ PPT created: {filepath}")
        
        return jsonify({
            "success": True,
            "file_path": filepath,
            "filename": filename
        })
    
    except Exception as e:
        print(f"❌ Error generating PPT: {str(e)}")
        return jsonify({"success": False, "error": str(e)})

@app.route("/download/<filename>", methods=["GET"])
def download_ppt(filename):
    """Download generated PPT file"""
    try:
        filepath = os.path.join(PPT_DIR, filename)
        if os.path.exists(filepath):
            return send_file(filepath, as_attachment=True, download_name=filename)
        return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    # For desktop app, open browser automatically
    if is_bundled:
        import threading
        def open_browser():
            import time
            time.sleep(1)  # Wait for server to start
            webbrowser.open("http://localhost:5000")
        
        thread = threading.Thread(target=open_browser, daemon=True)
        thread.start()
        print("🌐 Opening browser in 1 second...")
        app.run(host='127.0.0.1', port=5000, debug=False)
    else:
        app.run(debug=True)

