import os
import subprocess
import platform
import winreg
from pptx import Presentation # type: ignore
import webbrowser

GENERATED_DIR = "generated_code"
os.makedirs(GENERATED_DIR, exist_ok=True)

# App name mappings to actual executable names
APP_ALIASES = {
    "discord": "discord.exe",
    "file manager": "explorer.exe",
    "explorer": "explorer.exe",
    "files": "explorer.exe",
    "whatsapp": "whatsapp.exe",
    "spotify": "spotify.exe",
    "telegram": "telegram.exe",
    "chrome": "chrome.exe",
    "firefox": "firefox.exe",
    "edge": "msedge.exe",
    "word": "winword.exe",
    "excel": "excel.exe",
    "powerpoint": "powerpnt.exe",
    "photoshop": "photoshop.exe",
    "visual studio": "devenv.exe",
    "vscode": "code.exe",
    "vs code": "code.exe",
}

def open_app(app_name):
    """Opens ANY app on Windows - comprehensive search and execution"""
    app_name_lower = app_name.lower().strip()
    
    print(f"🔍 Attempting to open: {app_name_lower}")
    
    # Check if there's an alias for this app
    if app_name_lower in APP_ALIASES:
        actual_app = APP_ALIASES[app_name_lower]
        print(f"  Found alias: {app_name_lower} → {actual_app}")
        app_name_lower = actual_app
    
    # Method 1: Direct shell execution with verification
    try:
        print(f"  [1/8] Direct shell execution...")
        result = subprocess.run(app_name_lower, shell=True, capture_output=True, timeout=3)
        # Only return True if no error or process ran
        if result.returncode == 0:
            print(f"✓ Opened {app_name} via shell execution")
            return True
    except subprocess.TimeoutExpired:
        # Timeout might mean app is starting - consider it success
        print(f"✓ Opened {app_name} via shell execution (timeout=app starting)")
        return True
    except Exception as e:
        print(f"  [1/8] Failed: {e}")
    
    # Method 2: With .exe extension via shell
    try:
        print(f"  [2/8] Shell with .exe extension...")
        result = subprocess.run(app_name_lower + ".exe", shell=True, capture_output=True, timeout=3)
        if result.returncode == 0:
            print(f"✓ Opened {app_name}.exe via shell")
            return True
    except subprocess.TimeoutExpired:
        print(f"✓ Opened {app_name}.exe via shell (timeout=app starting)")
        return True
    except Exception as e:
        print(f"  [2/8] Failed")
    
    # Method 3: os.startfile (Windows native)
    try:
        print(f"  [3/8] os.startfile...")
        os.startfile(app_name_lower)
        print(f"✓ Opened {app_name} via os.startfile")
        return True
    except Exception as e:
        print(f"  [3/8] Failed")
    
    # Method 4: Search Program Files deeply
    print(f"  [4/8] Searching Program Files...")
    for prog_dir in [os.getenv('ProgramFiles'), os.getenv('ProgramFiles(x86)')]:
        if not prog_dir or not os.path.exists(prog_dir):
            continue
        
        try:
            # Search for any exe matching app name
            for root, dirs, files in os.walk(prog_dir):
                for file in files:
                    if file.lower().endswith('.exe'):
                        if app_name_lower in file.lower() or file.lower().startswith(app_name_lower.split()[0]):
                            exe_path = os.path.join(root, file)
                            try:
                                subprocess.Popen([exe_path])
                                print(f"✓ Found and opened: {exe_path}")
                                return True
                            except:
                                pass
        except:
            pass
    
    # Method 5: Search AppData Local
    print(f"  [5/8] Searching AppData...")
    try:
        appdata = os.getenv('LOCALAPPDATA')
        if appdata:
            for root, dirs, files in os.walk(appdata):
                for file in files:
                    if file.lower().endswith('.exe'):
                        if app_name_lower in file.lower():
                            exe_path = os.path.join(root, file)
                            try:
                                subprocess.Popen([exe_path])
                                print(f"✓ Found and opened: {exe_path}")
                                return True
                            except:
                                pass
    except:
        pass
    
    # Method 6: PowerShell Get-Command (for apps in PATH)
    print(f"  [6/8] PowerShell Get-Command...")
    try:
        ps_script = f"(Get-Command {app_name_lower} -ErrorAction SilentlyContinue).Source"
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=5
        )
        if result.returncode == 0 and result.stdout.strip():
            exe_path = result.stdout.strip()
            try:
                subprocess.Popen([exe_path])
                print(f"✓ Opened via PowerShell: {exe_path}")
                return True
            except:
                pass
    except:
        pass
    
    # Method 7: Windows registry search
    print(f"  [7/8] Registry search...")
    try:
        reg_paths = [
            r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall",
            r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\Uninstall"
        ]
        
        for reg_path in reg_paths:
            try:
                with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, reg_path) as key:
                    for i in range(winreg.QueryInfoKey(key)[0]):
                        try:
                            subkey_name = winreg.EnumKey(key, i)
                            with winreg.OpenKey(key, subkey_name) as subkey:
                                display_name = winreg.QueryValueEx(subkey, "DisplayName")[0]
                                
                                if app_name_lower in display_name.lower():
                                    try:
                                        uninstall_str = winreg.QueryValueEx(subkey, "UninstallString")[0]
                                        
                                        # Extract exe path
                                        if '"' in uninstall_str:
                                            exe_path = uninstall_str.split('"')[1]
                                        else:
                                            exe_path = uninstall_str.split()[0]
                                        
                                        if exe_path and os.path.exists(exe_path):
                                            subprocess.Popen([exe_path])
                                            print(f"✓ Opened from registry: {display_name}")
                                            return True
                                    except:
                                        pass
                        except:
                            pass
            except:
                pass
    except:
        pass
    
    # Method 8: UWP apps via PowerShell (Microsoft Store apps)
    print(f"  [8/8] UWP Store apps...")
    try:
        ps_script = f"""
$app = Get-AppxPackage -Name "*{app_name_lower}*" -ErrorAction SilentlyContinue | Select-Object -First 1
if ($app) {{
    explorer.exe shell:appsFolder\\$($app.PackageFullName)!
    exit 0
}}
exit 1
"""
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=5
        )
        if result.returncode == 0:
            print(f"✓ Opened UWP app: {app_name}")
            return True
    except:
        pass
    
    print(f"✗ Could not find or open: {app_name}")
    return False

def get_website_for_app(app_name):
    """Get official website/web version for an app"""
    app_lower = app_name.lower().strip()
    
    websites = {
        "whatsapp": "https://web.whatsapp.com/",
        "spotify": "https://open.spotify.com/",
        "telegram": "https://web.telegram.org/",
        "discord": "https://discord.com/",
        "slack": "https://app.slack.com/",
        "microsoft store": "https://www.microsoft.com/en-us/store",
        "store": "https://www.microsoft.com/en-us/store",
        "github": "https://github.com/",
        "gmail": "https://mail.google.com/",
    }
    
    return websites.get(app_lower, f"https://www.google.com/search?q={app_name}")

def execute_plan(plan):
    results = []

    for step in plan.get("actions", []):
        action = step.get("type")

        # =======================
        # OPEN APPLICATION
        # =======================
        if action == "open_application":
            name = step.get("name", "").strip()

            if open_app(name):
                results.append(f"✓ Opened {name}")
            else:
                # Fallback: Try to open website
                website = get_website_for_app(name)
                try:
                    webbrowser.open(website)
                    results.append(f"{name} app not found. Opened web version: {website}")
                except:
                    results.append(f"Could not open {name} or web version")

        # =======================
        # OPEN WEBSITE
        # =======================
        elif action == "open_website":
            url = step.get("url", "")
            try:
                webbrowser.open(url)
                results.append(f"✓ Opened website: {url}")
            except Exception as e:
                results.append(f"Could not open website: {e}")

        # =======================
        # WRITE FILE
        # =======================
        elif action == "write_file":
            path = os.path.join(GENERATED_DIR, step["file_name"])
            with open(path, "w", encoding="utf-8") as f:
                f.write(step["content"])
            if platform.system() == "Windows":
                os.startfile(path)
            else:
                subprocess.Popen(["xdg-open", path])
            results.append(f"✓ Wrote and opened {step['file_name']}")

        # =======================
        # CREATE PPT (REAL)
        # =======================
        elif action == "create_ppt":
            prs = Presentation()

            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = step["title"]

            for s in step["slides"]:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = s["title"]
                slide.placeholders[1].text = s["content"]

            path = os.path.join(GENERATED_DIR, step["title"] + ".pptx")
            prs.save(path)

            subprocess.Popen(path, shell=True)
            results.append("✓ Created and opened PowerPoint")

        # =======================
        # RESPOND
        # =======================
        elif action == "respond":
            results.append(step["message"])

    return " | ".join(results)
